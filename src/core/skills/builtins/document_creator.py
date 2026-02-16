"""
Built-in skill: document_creator — create PDF, Word, Excel, and PowerPoint documents.

Generates professional documents from structured content inputs.
Operates only within allowed paths (enforced by workspace sandboxing).
"""

from __future__ import annotations

import logging
import os
from pathlib import Path
from typing import Any, Dict, List

logger = logging.getLogger(__name__)

MANIFEST = {
    "name": "document_creator",
    "version": "1.0.0",
    "description": "Create PDF, Word, Excel, and PowerPoint documents",
    "risk": "MEDIUM",
    "permissions": ["file_write"],
    "inputs": [
        {"name": "format", "type": "string", "required": True,
         "description": "pdf|docx|xlsx|pptx"},
        {"name": "path", "type": "string", "required": True,
         "description": "Output file path (relative to workspace)"},
        {"name": "content", "type": "object", "required": True,
         "description": "Document content structure"},
    ],
}

DEFAULT_WORKSPACE = os.getenv("LANCELOT_WORKSPACE", "/home/lancelot/data")


def execute(context, inputs: Dict[str, Any]) -> Dict[str, Any]:
    """Create a document in the requested format."""
    fmt = inputs.get("format", "").lower()
    rel_path = inputs.get("path", "")
    content = inputs.get("content", {})
    workspace = inputs.get("workspace", DEFAULT_WORKSPACE)

    if not rel_path:
        raise ValueError("Missing required input: 'path'")

    if fmt not in ("pdf", "docx", "xlsx", "pptx"):
        raise ValueError(f"Unknown format: '{fmt}'. Must be pdf|docx|xlsx|pptx")

    # Ensure correct extension
    if not rel_path.lower().endswith(f".{fmt}"):
        rel_path = f"{rel_path}.{fmt}"

    full_path = _resolve_safe_path(workspace, rel_path)
    full_path.parent.mkdir(parents=True, exist_ok=True)

    if fmt == "pdf":
        return _create_pdf(full_path, content)
    elif fmt == "docx":
        return _create_docx(full_path, content)
    elif fmt == "xlsx":
        return _create_xlsx(full_path, content)
    elif fmt == "pptx":
        return _create_pptx(full_path, content)

    return {"status": "error", "error": f"Unhandled format: {fmt}"}


def _resolve_safe_path(workspace: str, rel_path: str) -> Path:
    ws = Path(workspace).resolve()
    target = (ws / rel_path).resolve()
    if not str(target).startswith(str(ws)):
        raise ValueError(f"Path traversal blocked: '{rel_path}' escapes workspace")
    return target


# ---------------------------------------------------------------------------
# PDF (reportlab)
# ---------------------------------------------------------------------------

def _create_pdf(path: Path, content: Dict[str, Any]) -> Dict[str, Any]:
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.units import inch
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
    from reportlab.lib import colors

    doc = SimpleDocTemplate(str(path), pagesize=letter,
                            leftMargin=inch, rightMargin=inch,
                            topMargin=inch, bottomMargin=inch)
    styles = getSampleStyleSheet()

    # Custom styles
    title_style = ParagraphStyle(
        "DocTitle", parent=styles["Title"], fontSize=20, spaceAfter=20, alignment=TA_CENTER,
    )
    heading_style = ParagraphStyle(
        "DocHeading", parent=styles["Heading2"], fontSize=14, spaceBefore=16, spaceAfter=8,
    )
    body_style = styles["BodyText"]
    body_style.spaceAfter = 8

    story: List[Any] = []

    # Title
    if content.get("title"):
        story.append(Paragraph(content["title"], title_style))
        story.append(Spacer(1, 12))

    # Subtitle
    if content.get("subtitle"):
        sub_style = ParagraphStyle("Sub", parent=styles["Normal"], fontSize=12, alignment=TA_CENTER, textColor=colors.grey)
        story.append(Paragraph(content["subtitle"], sub_style))
        story.append(Spacer(1, 24))

    # Sections
    for section in content.get("sections", []):
        if section.get("heading"):
            story.append(Paragraph(section["heading"], heading_style))
        for para in section.get("paragraphs", []):
            story.append(Paragraph(para, body_style))
        # Bullet points
        for bullet in section.get("bullets", []):
            bullet_style = ParagraphStyle("Bullet", parent=body_style, leftIndent=20, bulletIndent=10)
            story.append(Paragraph(f"\u2022 {bullet}", bullet_style))
        story.append(Spacer(1, 8))

    # Standalone paragraphs (simple mode)
    for para in content.get("paragraphs", []):
        story.append(Paragraph(para, body_style))

    # Tables
    for table_def in content.get("tables", []):
        headers = table_def.get("headers", [])
        rows = table_def.get("rows", [])
        data = [headers] + rows if headers else rows
        if data:
            t = Table(data)
            t.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#2d3748")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, 0), 10),
                ("BOTTOMPADDING", (0, 0), (-1, 0), 8),
                ("BACKGROUND", (0, 1), (-1, -1), colors.HexColor("#f7fafc")),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                ("FONTSIZE", (0, 1), (-1, -1), 9),
                ("TOPPADDING", (0, 1), (-1, -1), 4),
                ("BOTTOMPADDING", (0, 1), (-1, -1), 4),
            ]))
            story.append(t)
            story.append(Spacer(1, 12))

    doc.build(story)
    size = path.stat().st_size
    logger.info("document_creator: PDF created %s (%d bytes)", path, size)
    return {"status": "created", "path": str(path), "format": "pdf", "size_bytes": size}


# ---------------------------------------------------------------------------
# Word (python-docx)
# ---------------------------------------------------------------------------

def _create_docx(path: Path, content: Dict[str, Any]) -> Dict[str, Any]:
    from docx import Document
    from docx.shared import Pt, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # Title
    if content.get("title"):
        title_para = doc.add_heading(content["title"], level=0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Subtitle
    if content.get("subtitle"):
        sub = doc.add_paragraph()
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = sub.add_run(content["subtitle"])
        run.font.size = Pt(12)
        run.font.italic = True

    # Sections
    for section in content.get("sections", []):
        if section.get("heading"):
            doc.add_heading(section["heading"], level=2)
        for para in section.get("paragraphs", []):
            doc.add_paragraph(para)
        for bullet in section.get("bullets", []):
            doc.add_paragraph(bullet, style="List Bullet")

    # Standalone paragraphs
    for para in content.get("paragraphs", []):
        doc.add_paragraph(para)

    # Tables
    for table_def in content.get("tables", []):
        headers = table_def.get("headers", [])
        rows = table_def.get("rows", [])
        if headers or rows:
            cols = len(headers) if headers else len(rows[0]) if rows else 0
            total_rows = (1 if headers else 0) + len(rows)
            table = doc.add_table(rows=total_rows, cols=cols)
            table.style = "Table Grid"
            row_idx = 0
            if headers:
                for i, h in enumerate(headers):
                    cell = table.rows[0].cells[i]
                    cell.text = str(h)
                    for run in cell.paragraphs[0].runs:
                        run.bold = True
                row_idx = 1
            for row_data in rows:
                for i, val in enumerate(row_data):
                    table.rows[row_idx].cells[i].text = str(val)
                row_idx += 1

    doc.save(str(path))
    size = path.stat().st_size
    logger.info("document_creator: DOCX created %s (%d bytes)", path, size)
    return {"status": "created", "path": str(path), "format": "docx", "size_bytes": size}


# ---------------------------------------------------------------------------
# Excel (openpyxl)
# ---------------------------------------------------------------------------

def _create_xlsx(path: Path, content: Dict[str, Any]) -> Dict[str, Any]:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    wb = Workbook()

    sheets_data = content.get("sheets", [])
    if not sheets_data:
        # Simple mode: single sheet from top-level headers/rows
        sheets_data = [{
            "name": content.get("title", "Sheet1"),
            "headers": content.get("headers", []),
            "rows": content.get("rows", []),
        }]

    for idx, sheet_def in enumerate(sheets_data):
        ws = wb.active if idx == 0 else wb.create_sheet()
        ws.title = sheet_def.get("name", f"Sheet{idx + 1}")

        headers = sheet_def.get("headers", [])
        rows = sheet_def.get("rows", [])

        header_font = Font(bold=True, color="FFFFFF", size=11)
        header_fill = PatternFill(start_color="2D3748", end_color="2D3748", fill_type="solid")
        thin_border = Border(
            left=Side(style="thin"), right=Side(style="thin"),
            top=Side(style="thin"), bottom=Side(style="thin"),
        )

        if headers:
            for col, h in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col, value=h)
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = Alignment(horizontal="center")
                cell.border = thin_border

        start_row = 2 if headers else 1
        for row_idx, row_data in enumerate(rows, start_row):
            for col_idx, val in enumerate(row_data, 1):
                cell = ws.cell(row=row_idx, column=col_idx, value=val)
                cell.border = thin_border

        # Auto-width columns
        for col in ws.columns:
            max_len = 0
            col_letter = col[0].column_letter
            for cell in col:
                if cell.value:
                    max_len = max(max_len, len(str(cell.value)))
            ws.column_dimensions[col_letter].width = min(max_len + 4, 50)

    wb.save(str(path))
    size = path.stat().st_size
    logger.info("document_creator: XLSX created %s (%d bytes)", path, size)
    return {"status": "created", "path": str(path), "format": "xlsx", "size_bytes": size}


# ---------------------------------------------------------------------------
# PowerPoint (python-pptx)
# ---------------------------------------------------------------------------

def _create_pptx(path: Path, content: Dict[str, Any]) -> Dict[str, Any]:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()

    # Title slide
    if content.get("title"):
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide layout
        slide.shapes.title.text = content["title"]
        if content.get("subtitle") and slide.placeholders[1]:
            slide.placeholders[1].text = content["subtitle"]

    # Content slides from sections
    for section in content.get("sections", []):
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
        if section.get("heading"):
            slide.shapes.title.text = section["heading"]
        body = slide.placeholders[1].text_frame
        body.clear()
        for para_text in section.get("paragraphs", []):
            p = body.add_paragraph()
            p.text = para_text
            p.font.size = Pt(16)
        for bullet in section.get("bullets", []):
            p = body.add_paragraph()
            p.text = bullet
            p.font.size = Pt(14)
            p.level = 1

    # Simple paragraphs as individual slides
    for para in content.get("paragraphs", []):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = ""
        body = slide.placeholders[1].text_frame
        body.clear()
        p = body.add_paragraph()
        p.text = para
        p.font.size = Pt(18)

    prs.save(str(path))
    size = path.stat().st_size
    logger.info("document_creator: PPTX created %s (%d bytes)", path, size)
    return {"status": "created", "path": str(path), "format": "pptx", "size_bytes": size}
